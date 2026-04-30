import os
import re
import shutil
import uuid
from pathlib import Path
from fastapi import FastAPI, File, UploadFile, Form, HTTPException, BackgroundTasks
from fastapi.responses import FileResponse, HTMLResponse
from fastapi.staticfiles import StaticFiles
from bs4 import BeautifulSoup
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

app = FastAPI(title="SlideToEditable")

BASE_DIR = Path(__file__).parent.resolve()
STATIC_DIR = BASE_DIR / "static"
UPLOAD_DIR = BASE_DIR / "uploads"
OUTPUT_DIR = BASE_DIR / "outputs"

UPLOAD_DIR.mkdir(exist_ok=True)
OUTPUT_DIR.mkdir(exist_ok=True)

SYSTEM_CHROME_PATHS = [
    r"C:\Program Files\Google\Chrome\Application\chrome.exe",
    r"C:\Program Files (x86)\Google\Chrome\Application\chrome.exe",
    os.path.expanduser(r"~\AppData\Local\Google\Chrome\Application\chrome.exe"),
]

CHROME_PATH = None
for cp in SYSTEM_CHROME_PATHS:
    if os.path.isfile(cp):
        CHROME_PATH = cp
        break

BROWSER_LAUNCH_OPTIONS = {"headless": True, "args": ["--no-sandbox"]}
if CHROME_PATH:
    BROWSER_LAUNCH_OPTIONS["executable_path"] = CHROME_PATH


app.mount("/static", StaticFiles(directory=str(STATIC_DIR)), name="static")


@app.get("/")
async def root():
    index_path = STATIC_DIR / "index.html"
    if not index_path.exists():
        raise HTTPException(status_code=500, detail="index.html not found")
    return HTMLResponse(content=index_path.read_text(encoding="utf-8"))


def _detect_structure_type(page) -> str:
    """检测 HTML 的结构类型: 'slide' (class=slide) 或 'presentation' (class=presentation+card) 或 'unknown'"""
    has_slides = page.evaluate("document.querySelectorAll('.slide').length > 0")
    if has_slides:
        return "slide"
    has_presentations = page.evaluate("document.querySelectorAll('.presentation').length > 0")
    if has_presentations:
        return "presentation"
    return "unknown"


def _prepare_pdf_slide_structure(page):
    """处理 slide 结构的 HTML 用于 PDF 导出"""
    page.evaluate("""
        (() => {
            const slides = document.querySelectorAll('.slide');
            if (!slides.length) return;
            slides.forEach(el => {
                el.style.setProperty('display', 'flex', 'important');
                el.style.setProperty('visibility', 'visible', 'important');
                el.style.setProperty('opacity', '1', 'important');
                el.style.setProperty('width', '100%', 'important');
                el.style.setProperty('max-width', '100%', 'important');
                el.style.setProperty('overflow', 'hidden', 'important');
                el.style.setProperty('overflow-y', 'hidden', 'important');
                el.classList.add('active');
            });
            const pdfBox = document.createElement('div');
            pdfBox.id = 'pdf-container';
            pdfBox.style.setProperty('display', 'block', 'important');
            pdfBox.style.setProperty('width', '100%', 'important');
            slides.forEach(el => pdfBox.appendChild(el));
            document.body.innerHTML = '';
            document.body.appendChild(pdfBox);
            document.body.style.setProperty('overflow', 'visible', 'important');
            document.body.style.setProperty('height', 'auto', 'important');
            document.body.style.setProperty('display', 'block', 'important');
            document.body.style.setProperty('margin', '0', 'important');
            document.body.style.setProperty('padding', '0', 'important');
            const style = document.createElement('style');
            style.textContent = '.slide:not(:first-child) { page-break-before: always !important; }';
            document.head.appendChild(style);
        })();
    """)


def _prepare_pdf_presentation_structure(page):
    """处理 presentation+card 结构的 HTML 用于 PDF 导出：
       1. 模拟点击所有导航按钮，触发可能绑定的 JS 事件来加载内容
       2. 显示所有 presentation
       3. 展开所有 card 内的 content-block
       4. 隐藏导航按钮和 toggle 按钮等干扰元素
    """
    page.evaluate("""
        (() => {
            // 0. First, click all nav buttons to trigger any JS events that load content
            const navButtons = document.querySelectorAll('nav button, .nav-btn, [data-pres]');
            navButtons.forEach(btn => {
                btn.click();
                btn.dispatchEvent(new MouseEvent('click', {bubbles: true, cancelable: true}));
            });
            // Also click all toggle buttons to expand content
            document.querySelectorAll('.toggle-btn').forEach(btn => {
                btn.click();
                btn.dispatchEvent(new MouseEvent('click', {bubbles: true, cancelable: true}));
            });
        })();
    """)
    page.wait_for_timeout(500)

    page.evaluate("""
        (() => {
            // 1. Show all presentations
            document.querySelectorAll('.presentation').forEach(el => {
                el.style.setProperty('display', 'block', 'important');
                el.classList.add('active');
            });
            // 2. Force expand all content-blocks (theory, task, solution)
            document.querySelectorAll('.content-block').forEach(el => {
                el.classList.add('open');
                el.style.setProperty('max-height', 'none', 'important');
                el.style.setProperty('height', 'auto', 'important');
                el.style.setProperty('overflow', 'visible', 'important');
                el.style.setProperty('transition', 'none', 'important');
            });
            // 3. Hide nav
            const nav = document.querySelector('nav');
            if (nav) {
                nav.style.setProperty('display', 'none', 'important');
            }
            // 4. Hide toggle buttons
            document.querySelectorAll('.toggle-btn').forEach(el => {
                el.style.setProperty('display', 'none', 'important');
            });
            // 5. Add page-break between presentations
            const style = document.createElement('style');
            style.textContent = '.presentation { page-break-before: always !important; } .presentation:first-of-type { page-break-before: avoid !important; }';
            document.head.appendChild(style);
        })();
    """)


def _prepare_pptx_presentation_structure(page, pres_index: int):
    """处理 presentation+card 结构：只显示指定 presentation，隐藏其他"""
    page.evaluate(f"""
        (() => {{
            document.querySelectorAll('.presentation').forEach((el, idx) => {{
                if (idx === {pres_index}) {{
                    el.style.setProperty('display', 'block', 'important');
                    el.classList.add('active');
                }} else {{
                    el.style.setProperty('display', 'none', 'important');
                }}
            }});
            // Expand all content-blocks in visible presentation
            document.querySelectorAll('.content-block').forEach(el => {{
                el.classList.add('open');
                el.style.setProperty('max-height', 'none', 'important');
                el.style.setProperty('overflow', 'visible', 'important');
            }});
            // Hide toggle buttons
            document.querySelectorAll('.toggle-btn').forEach(el => {{
                el.style.setProperty('display', 'none', 'important');
            }});
            // Hide nav
            const nav = document.querySelector('nav');
            if (nav) nav.style.setProperty('display', 'none', 'important');
        }})();
    """)


def convert_to_pdf(html_path: str, output_path: str, lang: str):
    from playwright.sync_api import sync_playwright

    with sync_playwright() as p:
        browser = p.chromium.launch(**BROWSER_LAUNCH_OPTIONS)
        context = browser.new_context(
            viewport={"width": 1280, "height": 720},
            device_scale_factor=2,
        )
        page = context.new_page()

        file_url = "file:///" + html_path.replace("\\", "/")
        page.goto(file_url, wait_until="load")

        try:
            page.evaluate(f"""
                (() => {{
                    document.querySelectorAll('.content-ru, .content-en, .content-zh').forEach(el => {{
                        el.classList.remove('active');
                    }});
                    document.querySelectorAll('.content-{lang}').forEach(el => {{
                        el.classList.add('active');
                    }});
                    const langTextMap = {{'ru': ['ru','рус','рос'], 'en': ['en','eng','анг'], 'zh': ['zh','cn','chi','中文','кит']}};
                    const texts = langTextMap['{lang}'] || ['{lang}'];
                    document.querySelectorAll('.lang-btn').forEach(b => {{
                        const lower = b.textContent.toLowerCase().trim();
                        if (texts.some(t => lower.includes(t))) {{
                            b.dispatchEvent(new MouseEvent('click', {{bubbles: true}}));
                        }}
                    }});
                }})();
            """)
            page.wait_for_timeout(500)
        except Exception:
            pass

        hide_script = """
        (() => {
            const selectors = [
                '.nav-btn', '.nav-button', '.navigation',
                '.download-btn', '.download-button', '.download',
                '.progress-bar', '.progress',
                '.control-btn', '.controls',
                'button', '.btn',
                '.lang-btn', '.language-selector',
                '.toolbar', '.menu',
            ];
            selectors.forEach(sel => {
                document.querySelectorAll(sel).forEach(el => {
                    el.style.display = 'none';
                    el.style.visibility = 'hidden';
                    el.style.opacity = '0';
                    el.style.pointerEvents = 'none';
                });
            });
        })();
        """
        page.evaluate(hide_script)
        page.wait_for_timeout(300)

        struct_type = _detect_structure_type(page)

        if struct_type == "slide":
            _prepare_pdf_slide_structure(page)
        elif struct_type == "presentation":
            _prepare_pdf_presentation_structure(page)
        else:
            # Fallback: just render the page as-is
            pass

        page.wait_for_timeout(200)

        page.pdf(
            path=output_path,
            format="A4",
            print_background=True,
            margin={"top": "0.5in", "bottom": "0.5in", "left": "0.5in", "right": "0.5in"},
        )

        context.close()
        browser.close()


def extract_hex_color(style_value: str) -> str | None:
    if not style_value:
        return None
    match = re.search(r'#([0-9a-fA-F]{6}|[0-9a-fA-F]{3})', style_value)
    if match:
        code = match.group(1)
        if len(code) == 3:
            code = "".join(c * 2 for c in code)
        return "#" + code
    match = re.search(r'background(?:-color)?\s*:\s*rgb\s*\(\s*(\d+)\s*,\s*(\d+)\s*,\s*(\d+)\s*\)', style_value, re.IGNORECASE)
    if match:
        r, g, b = int(match.group(1)), int(match.group(2)), int(match.group(3))
        return f"#{r:02x}{g:02x}{b:02x}"
    return None


def _get_pptx_slide_elements(soup, struct_type: str, lang: str):
    """获取 PPTX 的幻灯片元素列表（每个元素将生成一页幻灯片）"""
    if struct_type == "slide":
        return soup.find_all(class_="slide")
    elif struct_type == "presentation":
        # For presentation structure, each card becomes a slide
        cards = soup.find_all(class_="card")
        if cards:
            return cards
        # Fallback: use presentation divs
        return soup.find_all(class_="presentation")
    return []


def _extract_structured_text(slide_elem, lang: str) -> list[dict]:
    """从幻灯片元素中提取结构化文本，保留标题/正文层级关系
    返回: [{"level": "title"|"subtitle"|"heading"|"body"|"label", "text": str}, ...]
    """
    items = []
    content_blocks = slide_elem.select(f".content-{lang}")
    if not content_blocks:
        content_blocks = slide_elem.select('[class*="content-"]')
    if not content_blocks:
        content_blocks = [slide_elem]

    for block in content_blocks:
        for el in block.find_all(["h1", "h2", "h3", "h4", "h5", "h6", "p", "li", "div", "span"]):
            text = el.get_text(strip=True)
            if not text or len(text) < 2:
                continue

            tag = el.name
            cls = " ".join(el.get("class", []))

            if tag in ("h1",):
                items.append({"level": "title", "text": text})
            elif tag in ("h2",):
                items.append({"level": "subtitle", "text": text})
            elif tag in ("h3", "h4", "h5", "h6"):
                items.append({"level": "heading", "text": text})
            elif "section-title" in cls:
                items.append({"level": "label", "text": text})
            elif "card-header" in cls:
                items.append({"level": "heading", "text": text})
            elif "solution" in cls:
                items.append({"level": "body", "text": text})
            elif tag == "li":
                items.append({"level": "body", "text": f"• {text}"})
            elif tag == "p":
                items.append({"level": "body", "text": text})
            elif tag == "div" and text and len(text) > 20:
                items.append({"level": "body", "text": text})
    return items


def _build_pptx_slide_from_text(prs, slide_items: list[dict], slide_title: str = ""):
    """根据结构化文本构建纯文字 PPTX 幻灯片"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout

    left = Inches(0.4)
    top = Inches(0.2)
    width = Inches(9.2)
    height = Inches(5.2)

    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    first = True
    for item in slide_items:
        level = item["level"]
        text = item["text"]

        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()

        p.text = text
        p.space_after = Pt(4)

        if level == "title":
            p.font.size = Pt(24)
            p.font.bold = True
            p.font.color.rgb = RGBColor(0x1e, 0x40, 0xaf)
            p.space_after = Pt(8)
        elif level == "subtitle":
            p.font.size = Pt(20)
            p.font.bold = True
            p.font.color.rgb = RGBColor(0x25, 0x63, 0xeb)
            p.space_after = Pt(6)
        elif level == "heading":
            p.font.size = Pt(16)
            p.font.bold = True
            p.font.color.rgb = RGBColor(0x0f, 0x17, 0x2a)
            p.space_after = Pt(4)
        elif level == "label":
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = RGBColor(0x1e, 0x40, 0xaf)
            p.space_before = Pt(6)
            p.space_after = Pt(2)
        else:
            p.font.size = Pt(12)
            p.font.color.rgb = RGBColor(0x0f, 0x17, 0x2a)

    return slide


def convert_to_pptx(html_path: str, output_path: str, lang: str):
    from playwright.sync_api import sync_playwright

    # First try extracting from static HTML via BeautifulSoup
    with open(html_path, "r", encoding="utf-8", errors="replace") as f:
        html_content = f.read()
    soup = BeautifulSoup(html_content, "lxml")

    # Detect structure type from soup
    soup_slides = soup.find_all(class_="slide")
    soup_cards = soup.find_all(class_="card")
    soup_presentations = soup.find_all(class_="presentation")

    # If static HTML already has the elements, use them directly (no Playwright needed)
    if soup_slides or soup_cards:
        struct_type = "slide" if soup_slides else "presentation"
        soup_elements = soup_slides if soup_slides else soup_cards

        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(5.625)

        for elem in soup_elements:
            items = _extract_structured_text(elem, lang)
            if items:
                _build_pptx_slide_from_text(prs, items)

        prs.save(output_path)
        return

    # If elements are dynamically rendered by JS, use Playwright to render and extract
    with sync_playwright() as p:
        browser = p.chromium.launch(**BROWSER_LAUNCH_OPTIONS)
        context = browser.new_context(
            viewport={"width": 1280, "height": 720},
            device_scale_factor=2,
        )
        page = context.new_page()

        file_url = "file:///" + html_path.replace("\\", "/")
        page.goto(file_url, wait_until="load")

        # Language switching
        try:
            page.evaluate(f"""
                (() => {{
                    document.querySelectorAll('.content-ru, .content-en, .content-zh').forEach(el => {{
                        el.classList.remove('active');
                    }});
                    document.querySelectorAll('.content-{lang}').forEach(el => {{
                        el.classList.add('active');
                    }});
                    const langTextMap = {{'ru': ['ru','рус','рос'], 'en': ['en','eng','анг'], 'zh': ['zh','cn','chi','中文','кит']}};
                    const texts = langTextMap['{lang}'] || ['{lang}'];
                    document.querySelectorAll('.lang-btn').forEach(b => {{
                        const lower = b.textContent.toLowerCase().trim();
                        if (texts.some(t => lower.includes(t))) {{
                            b.dispatchEvent(new MouseEvent('click', {{bubbles: true}}));
                        }}
                    }});
                }})();
            """)
            page.wait_for_timeout(500)
        except Exception:
            pass

        # Click all nav buttons and toggle buttons to expand content
        page.evaluate("""
            (() => {
                const navButtons = document.querySelectorAll('nav button, .nav-btn, [data-pres]');
                navButtons.forEach(btn => {
                    btn.click();
                    btn.dispatchEvent(new MouseEvent('click', {bubbles: true, cancelable: true}));
                });
                document.querySelectorAll('.toggle-btn').forEach(btn => {
                    btn.click();
                    btn.dispatchEvent(new MouseEvent('click', {bubbles: true, cancelable: true}));
                });
            })();
        """)
        page.wait_for_timeout(500)

        # Detect structure from rendered DOM
        struct_type = _detect_structure_type(page)

        if struct_type == "unknown":
            # Fallback: try to find card elements
            has_cards = page.evaluate("document.querySelectorAll('.card').length > 0")
            if has_cards:
                struct_type = "presentation"
            else:
                context.close()
                browser.close()
                raise ValueError("未在 HTML 中找到任何 class='slide' 或 class='presentation' 的元素")

        # Extract rendered HTML and parse with BeautifulSoup
        rendered_html = page.content()
        context.close()
        browser.close()

    # Parse the rendered HTML to extract text
    rendered_soup = BeautifulSoup(rendered_html, "lxml")
    soup_elements = _get_pptx_slide_elements(rendered_soup, struct_type, lang)

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)

    for elem in soup_elements:
        items = _extract_structured_text(elem, lang)
        if items:
            _build_pptx_slide_from_text(prs, items)

    prs.save(output_path)


def safe_unlink(path: Path):
    try:
        if path.exists():
            path.unlink()
    except Exception:
        pass


def cleanup_files(input_path: Path, output_path: Path):
    safe_unlink(input_path)
    safe_unlink(output_path)


@app.post("/convert")
def convert(
    background_tasks: BackgroundTasks,
    file: UploadFile = File(...),
    type: str = Form(...),
    lang: str = Form("ru"),
):
    if not (file.filename or "").lower().endswith((".html", ".htm")):
        raise HTTPException(status_code=400, detail="仅支持 .html 文件")

    if type not in ("pdf", "pptx"):
        raise HTTPException(status_code=400, detail="type 必须是 pdf 或 pptx")

    input_filename = f"{uuid.uuid4().hex}_{file.filename}"
    input_path = UPLOAD_DIR / input_filename

    ext = ".pdf" if type == "pdf" else ".pptx"
    output_filename = f"{uuid.uuid4().hex}{ext}"
    output_path = OUTPUT_DIR / output_filename

    try:
        content = file.file.read()
        input_path.write_bytes(content)
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"文件读取失败: {e}")
    finally:
        file.file.close()

    try:
        if type == "pdf":
            convert_to_pdf(str(input_path), str(output_path), lang)
        else:
            convert_to_pptx(str(input_path), str(output_path), lang)
    except ValueError as e:
        safe_unlink(input_path)
        safe_unlink(output_path)
        raise HTTPException(status_code=400, detail=str(e))
    except Exception as e:
        safe_unlink(input_path)
        safe_unlink(output_path)
        raise HTTPException(status_code=500, detail=f"转换失败: {e}")

    if not output_path.exists():
        safe_unlink(input_path)
        raise HTTPException(status_code=500, detail="转换失败：未生成输出文件")

    download_name = file.filename.rsplit(".", 1)[0] + ext

    background_tasks.add_task(cleanup_files, input_path, output_path)

    return FileResponse(
        path=str(output_path),
        filename=download_name,
        media_type="application/octet-stream",
    )


if __name__ == "__main__":
    import uvicorn
    uvicorn.run("app:app", host="127.0.0.1", port=8001)